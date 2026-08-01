"""Build the OATM project presentation (python-pptx, no Node/Canva dependency).

Import the resulting .pptx into Canva via "Import PPTX" for further editing.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- Palette
NAVY = RGBColor(0x12, 0x17, 0x3A)
CARD_NAVY = RGBColor(0x1E, 0x27, 0x61)
AMBER = RGBColor(0xFF, 0xA6, 0x30)
AMBER_DEEP = RGBColor(0xE0, 0x84, 0x1C)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0xF3, 0xF5, 0xFA)
DARKGRAY = RGBColor(0x2E, 0x33, 0x48)
MUTED = RGBColor(0x69, 0x72, 0x8C)
LINE_COLOR = RGBColor(0xDD, 0xE2, 0xF0)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
RED = RGBColor(0xC0, 0x39, 0x2B)

HEADER_FONT = "Cambria"
BODY_FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.7)
CONTENT_W = SLIDE_W - 2 * MARGIN

RESULTS = r"C:\Users\HP\OneDrive\Desktop\srsi\results"
SAMPLES = r"C:\Users\HP\OneDrive\Desktop\srsi\occluded_samples"
PRES_ASSETS = r"C:\Users\HP\OneDrive\Desktop\srsi\presentation"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ================================================================== Helpers
def add_slide(bg_color):
    slide = prs.slides.add_slide(BLANK)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    return slide


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _spPr(shape):
    return shape._element.spPr


def add_soft_shadow(shape, blur_pt=16, dist_pt=5, alpha_pct=38, direction=2700000, color_hex="12173A"):
    spPr = _spPr(shape)
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    effectLst = spPr.makeelement(qn('a:effectLst'), {})
    outerShdw = effectLst.makeelement(qn('a:outerShdw'), {
        'blurRad': str(int(Pt(blur_pt))),
        'dist': str(int(Pt(dist_pt))),
        'dir': str(direction),
        'rotWithShape': '0',
    })
    clr = outerShdw.makeelement(qn('a:srgbClr'), {'val': color_hex})
    alphaEl = clr.makeelement(qn('a:alpha'), {'val': str(int(alpha_pct * 1000))})
    clr.append(alphaEl)
    outerShdw.append(clr)
    effectLst.append(outerShdw)
    spPr.append(effectLst)


def set_fill_alpha(shape, alpha_pct):
    spPr = _spPr(shape)
    solidFill = spPr.find(qn('a:solidFill'))
    srgbClr = solidFill.find(qn('a:srgbClr'))
    for el in srgbClr.findall(qn('a:alpha')):
        srgbClr.remove(el)
    alphaEl = srgbClr.makeelement(qn('a:alpha'), {'val': str(int(alpha_pct * 1000))})
    srgbClr.append(alphaEl)


def add_overlay_rect(slide, left, top, width, height, color, alpha_pct):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    set_fill_alpha(rect, alpha_pct)
    return rect


def add_text(slide, left, top, width, height, text, size, color, bold=False,
             italic=False, font=BODY_FONT, align=PP_ALIGN.LEFT, anchor=None,
             line_spacing=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color
    return box


def add_kicker(slide, left, top, text, color=AMBER, size=13, align=PP_ALIGN.LEFT, width=Inches(8)):
    return add_text(slide, left, top, width, Inches(0.35), text.upper(), size, color,
                     bold=True, font=BODY_FONT, align=align)


def add_accent_rule(slide, left, top, width=Inches(0.55), color=AMBER, thickness_pt=4):
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(thickness_pt))
    rule.fill.solid()
    rule.fill.fore_color.rgb = color
    rule.line.fill.background()
    rule.shadow.inherit = False
    return rule


def add_rounded_box(slide, left, top, width, height, fill_color, radius=0.08, shadow=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    try:
        shape.adjustments[0] = radius
    except Exception:
        pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    if shadow:
        add_soft_shadow(shape)
    return shape


def framed_picture(slide, path, left, top, width, height=None, border_color=LINE_COLOR, shadow=True):
    if height is None:
        pic = slide.shapes.add_picture(path, left, top, width=width)
    else:
        pic = slide.shapes.add_picture(path, left, top, width=width, height=height)
    try:
        pic.line.color.rgb = border_color
        pic.line.width = Pt(0.75)
    except Exception:
        pass
    if shadow:
        add_soft_shadow(pic, blur_pt=14, dist_pt=4, alpha_pct=32)
    return pic


def add_badge_row(slide, left, top, width, number, text, size=14.5, row_h=Inches(0.9),
                   text_color=DARKGRAY, badge_color=NAVY, number_color=WHITE):
    badge_d = Inches(0.5)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top + (row_h - badge_d) // 2,
                                    badge_d, badge_d)
    badge.adjustments[0] = 0.3
    badge.fill.solid()
    badge.fill.fore_color.rgb = badge_color
    badge.line.fill.background()
    badge.shadow.inherit = False
    tf = badge.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number)
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.name = BODY_FONT
    run.font.color.rgb = number_color

    text_left = left + badge_d + Inches(0.28)
    text_width = width - badge_d - Inches(0.28)
    box = slide.shapes.add_textbox(text_left, top, text_width, row_h)
    tf2 = box.text_frame
    tf2.word_wrap = True
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = text
    run2.font.size = Pt(size)
    run2.font.name = BODY_FONT
    run2.font.color.rgb = text_color
    return box


def add_pill_label(slide, left, top, width, text, fill_color, text_color=WHITE, size=13):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.42))
    try:
        shape.adjustments[0] = 0.5
    except Exception:
        pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.name = BODY_FONT
    run.font.color.rgb = text_color
    return shape


def add_stat_ribbon(slide, left, top, width, height, stats, card_color=CARD_NAVY,
                     big_size=34, small_size=12.5, big_color=AMBER, small_color=ICE):
    card = add_rounded_box(slide, left, top, width, height, card_color, radius=0.1, shadow=True)
    n = len(stats)
    col_w = width // n
    for i, (big, small) in enumerate(stats):
        x = left + i * col_w
        add_text(slide, x, top + Inches(0.18), col_w, Inches(0.7), big, big_size, big_color,
                  bold=True, font=HEADER_FONT, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.15), top + Inches(0.18) + Inches(0.68), col_w - Inches(0.3),
                  height - Inches(0.95), small, small_size, small_color, align=PP_ALIGN.CENTER,
                  line_spacing=1.05)
        if i > 0:
            div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top + Inches(0.22), Pt(1),
                                          height - Inches(0.44))
            div.fill.solid()
            div.fill.fore_color.rgb = ICE
            div.line.fill.background()
            div.shadow.inherit = False
            set_fill_alpha(div, 35)
    return card


def style_chart_simple(chart, color_hex, show_data_labels=True, num_fmt="0%"):
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = show_data_labels
    if show_data_labels:
        dls = plot.data_labels
        dls.number_format = num_fmt
        dls.number_format_is_linked = False
        dls.font.size = Pt(12.5)
        dls.font.bold = True
        dls.font.color.rgb = NAVY
        dls.position = XL_LABEL_POSITION.OUTSIDE_END
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = RGBColor.from_string(color_hex)
    chart.category_axis.tick_labels.font.size = Pt(11.5)
    chart.category_axis.tick_labels.font.color.rgb = DARKGRAY
    chart.category_axis.format.line.color.rgb = LINE_COLOR
    chart.value_axis.visible = False
    chart.value_axis.has_major_gridlines = False
    chart.category_axis.has_major_gridlines = False


def style_dual_chart(chart, title, color_a=MUTED, color_b=AMBER):
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    run0 = chart.chart_title.text_frame.paragraphs[0].runs[0]
    run0.font.size = Pt(13.5)
    run0.font.bold = True
    run0.font.color.rgb = NAVY
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10.5)
    chart.plots[0].series[0].format.fill.solid()
    chart.plots[0].series[0].format.fill.fore_color.rgb = color_a
    chart.plots[0].series[1].format.fill.solid()
    chart.plots[0].series[1].format.fill.fore_color.rgb = color_b
    chart.category_axis.tick_labels.font.size = Pt(10.5)
    chart.category_axis.tick_labels.font.color.rgb = DARKGRAY
    chart.category_axis.format.line.color.rgb = LINE_COLOR
    chart.value_axis.tick_labels.font.size = Pt(9.5)
    chart.value_axis.has_major_gridlines = False
    chart.category_axis.has_major_gridlines = False


# ================================================================== Slide 1
# Title - full-bleed hero photo (your hook goes right before this slide)
s = add_slide(NAVY)
framed_picture(s, SAMPLES + r"\sample_003\1_previous_no_occlusion.jpg",
               0, 0, width=SLIDE_W, height=SLIDE_H, shadow=False)
add_overlay_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY, 40)
add_overlay_rect(s, 0, Inches(3.6), SLIDE_W, Inches(3.9), NAVY, 62)

add_kicker(s, Inches(1), Inches(0.55), "SRSI Research Project", color=AMBER, size=14)
add_accent_rule(s, Inches(1), Inches(4.15), width=Inches(0.7))
add_text(s, Inches(1), Inches(4.35), Inches(11.33), Inches(1.1),
         "Occlusion-Adaptive Temporal Memory", 40, WHITE, bold=True, font=HEADER_FONT)
add_text(s, Inches(1), Inches(5.35), Inches(10.33), Inches(0.6),
         "Giving self-driving cars short-term memory for hidden objects", 19, ICE, italic=True)
add_text(s, Inches(1), Inches(6.85), Inches(10.33), Inches(0.4),
         "nuScenes CAM_FRONT — an autonomous-vehicle perception research project", 12.5, MUTED)
add_notes(s, "This follows straight after your hook - it's the reveal of what the project actually is. "
             "Introduce the project over the real dashcam frame: this is about helping self-driving car "
             "cameras keep track of road users even when something briefly blocks the view.")

# ================================================================== Slide 2
# The Problem - 3-frame filmstrip of the running example
s = add_slide(WHITE)
add_kicker(s, MARGIN, Inches(0.55), "The Problem in Autonomous Driving")
add_text(s, MARGIN, Inches(0.95), Inches(11.9), Inches(0.7),
         "Single-frame detectors forget what they can't see", 26, NAVY, bold=True, font=HEADER_FONT)
add_accent_rule(s, MARGIN, Inches(1.62))

strip_w = Inches(3.75)
strip_h = strip_w * 9 // 16
gap = Inches(0.35)
xs = [MARGIN, MARGIN + strip_w + gap, MARGIN + 2 * (strip_w + gap)]
strip_y = Inches(2.15)
frames = [
    ("1_previous_no_occlusion.jpg", "BEFORE", "Car clearly visible", NAVY),
    ("3_full_occlusion.jpg", "OCCLUDED", "Fully hidden from view", AMBER_DEEP),
    ("5_full_appearance.jpg", "AFTER", "Reappears down the road", NAVY),
]
for x, (fname, tag, caption, tag_color) in zip(xs, frames):
    framed_picture(s, SAMPLES + rf"\sample_003\{fname}", x, strip_y, width=strip_w, height=strip_h)
    add_pill_label(s, x + Inches(0.15), strip_y + strip_h - Inches(0.32), Inches(1.6), tag, tag_color, size=11)
    add_text(s, x, strip_y + strip_h + Inches(0.12), strip_w, Inches(0.4), caption, 13, MUTED,
              align=PP_ALIGN.CENTER)

for x in xs[:-1]:
    add_text(s, x + strip_w, strip_y + strip_h // 2 - Inches(0.2), gap, Inches(0.4), "→", 22,
              MUTED, bold=True, align=PP_ALIGN.CENTER)

add_text(s, MARGIN, Inches(5.55), CONTENT_W, Inches(0.55),
         "A detector that only looks at the middle frame has no idea a car is even there.",
         18, DARKGRAY, bold=True, align=PP_ALIGN.CENTER)
add_text(s, MARGIN, Inches(6.15), CONTENT_W, Inches(0.5),
         "“Like watching someone walk behind a parked bus — you don't assume they vanished.”",
         14, MUTED, italic=True, align=PP_ALIGN.CENTER)
add_notes(s, "Walk through this one real example frame by frame: the car is visible, then fully hidden "
             "behind other traffic, then reappears further down the road. Point out that most detectors "
             "process the middle frame in total isolation - to them, the car simply doesn't exist for "
             "that instant. That's the core problem this whole project is about.")

# ================================================================== Slide 3
# The Dataset - nuScenes, front camera only
s = add_slide(WHITE)
add_kicker(s, MARGIN, Inches(0.55), "The Dataset")
add_text(s, MARGIN, Inches(0.95), Inches(6.4), Inches(0.9),
         "nuScenes — Front Camera Only", 25, NAVY, bold=True, font=HEADER_FONT)
add_accent_rule(s, MARGIN, Inches(1.62))

dataset_rows = [
    "nuScenes: a large-scale AV dataset with synchronized cameras, LiDAR, radar, and 3D box "
    "annotations that include a per-object visibility rating",
    "This project uses only CAM_FRONT — one monocular camera, no LiDAR or other views feeding "
    "detection itself",
    "Front-only keeps the pipeline laptop-scale and isolates the occlusion-memory problem, "
    "instead of mixing in multi-camera fusion",
    "LiDAR and the 3D annotations aren't wasted, though — they come back as independent ground "
    "truth for evaluating OATM (next: Evaluation Methodology)",
]
y = Inches(2.1)
for i, text in enumerate(dataset_rows, 1):
    add_badge_row(s, MARGIN, y, Inches(6.5), i, text, size=13.5, row_h=Inches(1.15))
    y += Inches(1.25)

pic_left, pic_top, pic_w = Inches(7.5), Inches(2.15), Inches(5.15)
pic_h = pic_w * 9 // 16
framed_picture(s, SAMPLES + r"\sample_003\1_previous_no_occlusion.jpg", pic_left, pic_top, width=pic_w, height=pic_h)
add_text(s, pic_left, pic_top + pic_h + Inches(0.15), pic_w, Inches(0.5),
         "CAM_FRONT — the only view this project detects from", 12.5, MUTED, italic=True,
         align=PP_ALIGN.CENTER)
add_notes(s, "Be upfront about scope before going further: this is nuScenes, and specifically only the "
             "forward-facing camera - no LiDAR, no side/rear cameras, no sensor fusion feeding the "
             "detector. That's a deliberate simplification to keep the study laptop-scale and focused "
             "purely on the occlusion-memory question. The payoff comes later - nuScenes' LiDAR and 3D "
             "annotations aren't part of detection, but they become the independent ground truth used "
             "to evaluate OATM, which is the next thing after the design slide.")

# ================================================================== Slide 4
# Research Question & Hypothesis
s = add_slide(NAVY)
add_kicker(s, Inches(0), Inches(1.55), "The Core Research Question", align=PP_ALIGN.CENTER, width=SLIDE_W)
add_accent_rule(s, SLIDE_W // 2 - Inches(0.35), Inches(2.0), width=Inches(0.7))
add_text(s, Inches(1.3), Inches(2.35), Inches(10.73), Inches(2.1),
         "Can short-term visual memory help a camera keep track of temporarily hidden "
         "road users — without inventing false ones?", 29, WHITE, bold=True,
         font=HEADER_FONT, align=PP_ALIGN.CENTER, line_spacing=1.15)

card = add_rounded_box(s, Inches(1.9), Inches(4.85), Inches(9.53), Inches(1.75), CARD_NAVY,
                        radius=0.12, shadow=True)
tf = card.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.4)
tf.margin_right = Inches(0.4)
tf.margin_top = Inches(0.22)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Hypothesis"
run.font.size = Pt(13)
run.font.bold = True
run.font.color.rgb = AMBER
run.font.name = BODY_FONT
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(6)
run2 = p2.add_run()
run2.text = ("Occlusion-aware temporal memory preserves more valid tracks through short "
             "occlusions than single-frame detection or conventional tracking — without "
             "creating an unacceptable number of false, ghost tracks.")
run2.font.size = Pt(15.5)
run2.font.italic = True
run2.font.color.rgb = ICE
run2.font.name = BODY_FONT
add_notes(s, "State the research question plainly and pause on it. This is the question the whole "
             "project is trying to answer, and the hypothesis is what we expect but haven't yet proven.")

# ================================================================== Slide 5
# Evidence 1 - Baseline (method + charts + stats, merged onto one slide)
s = add_slide(WHITE)
add_kicker(s, MARGIN, Inches(0.45), "Evidence 1 · Baseline")
add_text(s, MARGIN, Inches(0.82), Inches(11.9), Inches(0.6),
         "Detection Collapses at Full Occlusion", 25, NAVY, bold=True, font=HEADER_FONT)
add_accent_rule(s, MARGIN, Inches(1.42))
add_text(s, MARGIN, Inches(1.55), CONTENT_W, Inches(0.4),
         "Pretrained YOLO nano, CPU-only, on 8 real nuScenes occlusion sequences — tracked "
         "through 5 stages, from clearly visible to fully hidden and back.", 13.5, DARKGRAY)

stages = ["No\nOcclusion", "First Partial\nOcclusion", "Full\nOcclusion",
          "First Partial\nAppearance", "Full\nAppearance"]
detection_rates = [1.0, 1.0, 0.125, 0.5, 1.0]
mean_conf = [0.774, 0.587, 0.029, 0.425, 0.845]

cd1 = CategoryChartData()
cd1.categories = stages
cd1.add_series("Detection rate", detection_rates)
gframe1 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, MARGIN, Inches(1.95),
                              Inches(5.85), Inches(2.5), cd1)
chart1 = gframe1.chart
chart1.has_title = True
chart1.chart_title.text_frame.text = "Detection rate by stage"
chart1.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(13.5)
chart1.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
chart1.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = NAVY
style_chart_simple(chart1, "1E2761", num_fmt="0%")

cd2 = CategoryChartData()
cd2.categories = stages
cd2.add_series("Mean confidence", mean_conf)
gframe2 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(6.85), Inches(1.95),
                              Inches(5.85), Inches(2.5), cd2)
chart2 = gframe2.chart
chart2.has_title = True
chart2.chart_title.text_frame.text = "Mean confidence by stage (0 = undetected)"
chart2.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(13.5)
chart2.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
chart2.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = NAVY
style_chart_simple(chart2, "FFA630", num_fmt="0.00")

img_left, img_top, img_w = MARGIN, Inches(4.6), Inches(3.3)
img_h = img_w * 9 // 16
framed_picture(s, PRES_ASSETS + r"\assets_yolo_only_crop_sample_003.jpg",
               img_left, img_top, width=img_w, height=img_h)
add_text(s, img_left, img_top + img_h + Inches(0.06), img_w, Inches(0.6),
          "Full occlusion, sample_003 — YOLO draws zero boxes on the target. "
          "Not misdetected — simply absent.", 10.5, MUTED, italic=True, align=PP_ALIGN.CENTER)

ribbon_left = img_left + img_w + Inches(0.3)
ribbon_w = CONTENT_W - img_w - Inches(0.3)
add_stat_ribbon(s, ribbon_left, Inches(4.6), ribbon_w, Inches(1.3), [
    ("100% → 12%", "detection rate at full occlusion"),
    ("0.77 → 0.03", "mean confidence at full occlusion"),
    ("100%", "correct class whenever detected"),
], big_size=22, small_size=10.5)
add_text(s, ribbon_left, Inches(6.0), ribbon_w, Inches(0.4),
          "Partial-occlusion stages have small sample sizes (n=2-3) — interpret with caution.",
          11, MUTED, italic=True, align=PP_ALIGN.LEFT)
add_notes(s, "Baseline setup: a plain, off-the-shelf YOLO detector, no memory at all. Point at the "
             "photo - that's the actual failure, a real frame where the detector produced zero boxes "
             "for the hidden target, not a simulated one. Detection rate stays perfect through partial "
             "occlusion then falls off a cliff at full occlusion. Confidence erodes before detection "
             "technically fails, and recovers once visible again.")

# ================================================================== Slide 6
# Evidence 2 - Naive memory (best/worst + stats, merged onto one slide)
s = add_slide(WHITE)
add_kicker(s, MARGIN, Inches(0.45), "Evidence 2 · Naive Memory")
add_text(s, MARGIN, Inches(0.82), Inches(11.9), Inches(0.6),
         "The Simplest Fix: Freeze the Last Box", 25, NAVY, bold=True, font=HEADER_FONT)
add_accent_rule(s, MARGIN, Inches(1.42))
add_text(s, MARGIN, Inches(1.55), CONTENT_W, Inches(0.35),
         "When YOLO loses the target, keep drawing its last known box instead of letting it vanish.",
         13.5, DARKGRAY)

img_w = Inches(5.2)
img_h = img_w * 9 // 16
img_gap = Inches(0.3)
total_w = img_w * 2 + img_gap
start_x = (SLIDE_W - total_w) // 2
right_x = start_x + img_w + img_gap
label_y = Inches(2.0)
img_y = Inches(2.42)

add_pill_label(s, start_x, label_y, Inches(2.2), "BEST CASE", GREEN, size=12)
add_pill_label(s, right_x, label_y, Inches(2.3), "WORST CASE", RED, size=12)

framed_picture(s, RESULTS + r"\last_seen_memory\comparisons\sample_003\comparison_2_reappearance.jpg",
               start_x, img_y, width=img_w, height=img_h)
framed_picture(s, RESULTS + r"\last_seen_memory\comparisons\sample_006\comparison_2_reappearance.jpg",
               right_x, img_y, width=img_w, height=img_h)

cap_y = img_y + img_h + Inches(0.08)
add_text(s, start_x, cap_y, img_w, Inches(0.4),
          "sample_003 — reappears next to a bus (83 px, IoU 0.34 — boxes barely overlap by a third)",
          11.5, MUTED, italic=True, align=PP_ALIGN.CENTER)
add_text(s, right_x, cap_y, img_w, Inches(0.4),
          "sample_006 — crossing an intersection (589 px, IoU 0.00)", 11.5, MUTED, italic=True,
          align=PP_ALIGN.CENTER)
add_text(s, MARGIN, cap_y + Inches(0.42), CONTENT_W, Inches(0.3),
          "Also tested: sample_001 (127px), sample_004 (423px), sample_005 (432px) — same pattern "
          "across all 5.", 11, MUTED, italic=True, align=PP_ALIGN.CENTER)

add_stat_ribbon(s, MARGIN, Inches(6.15), CONTENT_W, Inches(1.2), [
    ("331 px", "mean center error on reappearance"),
    ("0.085", "mean IoU (median: 0.000)"),
    ("3 / 5", "samples had zero box overlap"),
], big_size=26, small_size=11)
add_notes(s, "Introduce the naive fix: just freeze the last box while the target is hidden. Orange is "
             "the frozen memory box, green is the real detection once the object reappears. Best case "
             "they nearly line up; worst case the frozen box is nowhere near the real car. Not cherry-"
             "picked - all 5 valid samples follow the same pattern, mean error 331px, 3 of 5 with zero "
             "overlap at all.")

# ================================================================== Slide 7
# Evidence 3 - SORT motion prediction (results + caveats, merged onto one slide)
s = add_slide(WHITE)
add_kicker(s, MARGIN, Inches(0.45), "Evidence 3 · SORT Motion Prediction")
add_text(s, MARGIN, Inches(0.82), Inches(11.9), Inches(0.6),
         "Motion Prediction Helps — With Caveats", 24, NAVY, bold=True, font=HEADER_FONT)
add_accent_rule(s, MARGIN, Inches(1.42))
add_text(s, MARGIN, Inches(1.55), CONTENT_W, Inches(0.35),
         "A mini SORT-style tracker predicting velocity instead of freezing position — "
         "tested on 9 real car tracks.", 13.5, DARKGRAY)

gaps = ["1", "2", "3", "5", "8", "10"]
static_err = [12.778, 16.152, 22.286, 34.21, 51.616, 64.485]
sort_err = [3.472, 3.3, 4.588, 6.503, 14.295, 31.329]
static_iou = [0.887, 0.86, 0.817, 0.737, 0.628, 0.575]
sort_iou = [0.921, 0.905, 0.882, 0.838, 0.702, 0.57]

cd1 = CategoryChartData()
cd1.categories = gaps
cd1.add_series("Static memory", static_err)
cd1.add_series("SORT motion", sort_err)
gframe1 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, MARGIN, Inches(2.0),
                              Inches(5.85), Inches(2.55), cd1)
style_dual_chart(gframe1.chart, "Center error by gap length (px)")

cd2 = CategoryChartData()
cd2.categories = gaps
cd2.add_series("Static memory", static_iou)
cd2.add_series("SORT motion", sort_iou)
gframe2 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(6.85), Inches(2.0),
                              Inches(5.85), Inches(2.55), cd2)
style_dual_chart(gframe2.chart, "Mean IoU by gap length")

add_stat_ribbon(s, MARGIN, Inches(4.7), CONTENT_W, Inches(1.2), [
    ("46 → 17 px", "mean center error, static → SORT"),
    ("0.673 → 0.720", "mean IoU, static → SORT"),
    ("9 tracks", "cars only, across 3 clips"),
], big_size=23, small_size=10.5)

catch_card = add_rounded_box(s, MARGIN, Inches(6.0), CONTENT_W, Inches(1.35), SLATE, radius=0.06, shadow=True)
tf = catch_card.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_right = Inches(0.3)
tf.margin_top = Inches(0.16)
p = tf.paragraphs[0]
run = p.add_run()
run.text = "THE CATCH"
run.font.size = Pt(11.5)
run.font.bold = True
run.font.color.rgb = AMBER_DEEP
run.font.name = BODY_FONT
p2 = tf.add_paragraph()
p2.space_before = Pt(3)
run2 = p2.add_run()
run2.text = ("Almost the entire win comes from 2 of 9 tracks — on the dominant clip, static memory "
             "ties SORT. The reference boxes were Kalman-corrected (not raw YOLO), and long gaps "
             "bypass the tracker's own expiry rule. A real demonstration of motion extrapolation — "
             "not yet an independent validation of SORT.")
run2.font.size = Pt(12.5)
run2.font.color.rgb = DARKGRAY
run2.font.name = BODY_FONT
add_notes(s, "Motion prediction beats a frozen box clearly - 46px down to 17px, IoU up from 0.673 to "
             "0.720. But say the catch out loud: almost the whole win comes from 2 of 9 tracks, the "
             "reference boxes aren't fully independent, and long gaps bypass the tracker's own expiry "
             "rule. A real result, not yet a fully validated one.")

# ================================================================== Slide 8
# The Diagnosis + The Gap
s = add_slide(NAVY)
add_kicker(s, Inches(0), Inches(0.6), "The Diagnosis", align=PP_ALIGN.CENTER, width=SLIDE_W)
add_text(s, Inches(1.3), Inches(1.0), Inches(10.73), Inches(1.3),
          "A static memory box goes stale because it ignores motion.", 27, WHITE, bold=True,
          font=HEADER_FONT, align=PP_ALIGN.CENTER, line_spacing=1.12)
add_text(s, Inches(1.7), Inches(2.35), Inches(9.93), Inches(0.85),
          "Real objects keep moving while hidden. A memory system that only remembers where "
          "something was will drift further wrong the longer it stays hidden.", 14, ICE,
          align=PP_ALIGN.CENTER, line_spacing=1.2)

divider = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.67), Inches(3.35), Inches(4), Pt(1.5))
divider.fill.solid()
divider.fill.fore_color.rgb = AMBER
divider.line.fill.background()
divider.shadow.inherit = False
set_fill_alpha(divider, 55)

add_kicker(s, Inches(0), Inches(3.6), "The Gap in Existing Tools", align=PP_ALIGN.CENTER, width=SLIDE_W)

gap_rows = [
    ("Single-frame detection", "no memory at all — occluded means gone"),
    ("SORT", "cuts error 46→17px on smooth tracks — but N=9 cars, circular reference"),
    ("ByteTrack", "still needs some visible evidence to link to"),
    ("Fixed temporal memory", "blind to motion and its own ghost-track risk"),
]
col_w = Inches(5.95)
col_gap = Inches(0.35)
gx = [Inches(1.0), Inches(1.0) + col_w + col_gap]
gy = [Inches(4.15), Inches(5.15)]
positions = [(gx[0], gy[0]), (gx[1], gy[0]), (gx[0], gy[1]), (gx[1], gy[1])]
for (label, weakness), (x, y) in zip(gap_rows, positions):
    card = add_rounded_box(s, x, y, col_w, Inches(0.85), CARD_NAVY, radius=0.15)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label
    run.font.size = Pt(13.5)
    run.font.bold = True
    run.font.name = BODY_FONT
    run.font.color.rgb = AMBER
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = weakness
    run2.font.size = Pt(12.5)
    run2.font.name = BODY_FONT
    run2.font.color.rgb = ICE

add_text(s, Inches(1.7), Inches(6.2), Inches(9.93), Inches(0.5),
          "None decide hidden vs. lost vs. gone, or decay confidence accordingly — "
          "that's the gap OATM is built to fill.", 14.5, AMBER, bold=True, italic=True,
          align=PP_ALIGN.CENTER)
add_notes(s, "This is the pivot of the whole talk. First state the diagnosis: naive memory fails "
             "because it ignores motion, and even real motion prediction (as just shown) only helps "
             "partially. Then name the gap in every existing lightweight method - none of them combine "
             "occlusion-causality reasoning, motion prediction, and confidence decay. That's the bridge "
             "into the OATM design on the next slide.")

# ================================================================== Slide 9
# OATM Design
s = add_slide(WHITE)
add_kicker(s, MARGIN, Inches(0.5), "The Proposed Fix")
add_text(s, MARGIN, Inches(0.87), Inches(9), Inches(0.7),
          "OATM: Occlusion-Adaptive Temporal Memory", 24, NAVY, bold=True, font=HEADER_FONT)
add_accent_rule(s, MARGIN, Inches(1.55))

rows = [
    "Appearance memory: stores the last clear crop of the object, so it can be re-identified on reappearance",
    "Motion memory: a Kalman-style predict → observe → correct loop — the same mechanism already prototyped in Evidence 3, fixed to use an independent reference and real track-expiry",
    "Occlusion classifier: VISIBLE / OCCLUDED / LOST / EXITED — decides how much to trust the motion prediction right now",
    "Confidence decays smoothly, not on a fixed countdown; anti-ghost termination retires a track before it becomes a phantom",
]
y = Inches(1.95)
for i, text in enumerate(rows, 1):
    add_badge_row(s, MARGIN, y, Inches(6.6), i, text, size=13, row_h=Inches(1.25))
    y += Inches(1.35)

# state-flow diagram
box_w, box_h = Inches(2.5), Inches(0.72)
diag_x = Inches(8.4)
states = [("VISIBLE", NAVY), ("OCCLUDED", AMBER_DEEP), ("LOST / EXITED", MUTED)]
ys = [Inches(2.15), Inches(3.55), Inches(4.95)]
centers = []
for (label, color), yy in zip(states, ys):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, diag_x, yy, box_w, box_h)
    shp.adjustments[0] = 0.18
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    add_soft_shadow(shp, blur_pt=10, dist_pt=3, alpha_pct=28)
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = BODY_FONT
    centers.append((diag_x + box_w // 2, yy, yy + box_h))

for i in range(len(centers) - 1):
    cx, _, y_bottom = centers[i]
    _, y_top_next, _ = centers[i + 1]
    conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, cx, y_bottom, cx, y_top_next)
    conn.line.color.rgb = MUTED
    conn.line.width = Pt(2)

add_text(s, diag_x, Inches(5.85), box_w, Inches(0.9),
          "reappearance → back to VISIBLE", 11, MUTED, italic=True, align=PP_ALIGN.CENTER)
add_notes(s, "Introduce OATM's core mechanism, and be specific about the motion half since that's the "
             "part judges will probe: it's not a vague promise, it's the same predict-observe-correct "
             "Kalman loop already built and tested in Evidence 3 - just corrected to use an independent "
             "reference box and to actually enforce track-expiry, the two flaws that experiment exposed. "
             "Appearance memory (last clear crop) handles re-identification; motion memory (position + "
             "velocity) handles where to look. The occlusion classifier decides how much to trust the "
             "motion prediction at any moment, and confidence decays smoothly rather than on a fixed "
             "timer. Be clear this is the proposed design - it has not been implemented or benchmarked "
             "yet, that's the next phase.")

# ================================================================== Slide 10
# Evaluation Methodology - LiDAR + 3D annotations as independent ground truth
s = add_slide(WHITE)
add_kicker(s, MARGIN, Inches(0.5), "Evaluation Methodology")
add_text(s, MARGIN, Inches(0.87), Inches(11.5), Inches(0.7),
          "Ground Truth From LiDAR, Not the Camera Itself", 23, NAVY, bold=True, font=HEADER_FONT)
add_accent_rule(s, MARGIN, Inches(1.5))

eval_rows = [
    "Use nuScenes' 3D box annotations + LiDAR point cloud as ground truth, projected into "
    "CAM_FRONT — never derived from the front camera's own detections",
    "Run OATM exactly as designed: CAM_FRONT only, genuinely blind to the object during its "
    "occluded frames",
    "Score OATM's predicted position during occlusion against that projected LiDAR / 3D ground "
    "truth, not against another camera-based tracker",
    "This directly fixes the circularity in Evidence 3, where the reference boxes came from the "
    "same Kalman model being tested",
]
y = Inches(1.95)
for i, text in enumerate(eval_rows, 1):
    add_badge_row(s, MARGIN, y, CONTENT_W, i, text, size=14.5, row_h=Inches(1.05))
    y += Inches(1.15)

callout = add_rounded_box(s, MARGIN, Inches(6.55), CONTENT_W, Inches(0.7), CARD_NAVY, radius=0.12, shadow=True)
tf = callout.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_right = Inches(0.3)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "The front camera can't see through the occlusion — but LiDAR still can. That's what makes this an independent check."
run.font.size = Pt(14)
run.font.italic = True
run.font.bold = True
run.font.color.rgb = AMBER
run.font.name = BODY_FONT
add_notes(s, "This is the answer to 'how will you actually know OATM is right.' The key move: ground "
             "truth comes from LiDAR and the 3D box annotations, not from another camera-derived "
             "tracker - so it isn't circular the way the SORT experiment's evaluation was. The front "
             "camera is blind during occlusion by design; LiDAR isn't, so it gives an independent check "
             "on whether OATM's predicted position is actually correct, not just self-consistent.")

# ================================================================== Slide 11
# How OATM Compares (table)
s = add_slide(WHITE)
add_kicker(s, MARGIN, Inches(0.5), "How OATM Compares")
add_text(s, MARGIN, Inches(0.87), Inches(9), Inches(0.6),
          "A Design Comparison", 24, NAVY, bold=True, font=HEADER_FONT)
add_accent_rule(s, MARGIN, Inches(1.5))
add_text(s, MARGIN, Inches(1.6), CONTENT_W, Inches(0.4),
          "OATM has not been implemented or measured yet — this is a design-level comparison.",
          13, MUTED, italic=True)

table_data = [
    ("Method", "How it works", "Main weakness"),
    ("Single-frame detection", "Looks only at the current image", "Hidden object instantly disappears"),
    ("SORT", "Predicts motion — measured 46→17px error on 9 tracks", "Tiny sample; reference partly circular; expiry bypassed"),
    ("ByteTrack", "Links strong and weak detections", "Still needs some visible evidence"),
    ("Fixed temporal memory", "Keeps missing objects for N frames", "May forget early or keep ghosts too long"),
    ("OATM (proposed)", "Appearance + motion + occlusion-aware decay", "More complex; needs careful tuning"),
]
n_rows, n_cols = len(table_data), 3
tbl_top = Inches(2.15)
tbl_h = Inches(4.55)
card = add_rounded_box(s, MARGIN - Inches(0.12), tbl_top - Inches(0.12), CONTENT_W + Inches(0.24),
                        tbl_h + Inches(0.24), WHITE, radius=0.03, shadow=True)
tbl_shape = s.shapes.add_table(n_rows, n_cols, MARGIN, tbl_top, CONTENT_W, tbl_h)
table = tbl_shape.table
table.columns[0].width = Inches(2.9)
table.columns[1].width = Inches(5.3)
table.columns[2].width = Inches(3.73)

for r, row in enumerate(table_data):
    for c, val in enumerate(row):
        cell = table.cell(r, c)
        cell.text = val
        cell.margin_left = Inches(0.18)
        cell.margin_right = Inches(0.15)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = cell.text_frame.paragraphs[0]
        run = para.runs[0]
        run.font.name = BODY_FONT
        if r == 0:
            run.font.bold = True
            run.font.size = Pt(13.5)
            run.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
        else:
            run.font.size = Pt(12.5)
            run.font.color.rgb = DARKGRAY
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0) if r == n_rows - 1 else \
                (SLATE if r % 2 == 0 else WHITE)
            if r == n_rows - 1:
                run.font.bold = True
                run.font.color.rgb = NAVY
add_notes(s, "Position OATM against the alternatives already used in the field. The point isn't that "
             "the others are bad - it's that none of them explicitly reason about *why* an object "
             "disappeared, which is what OATM adds. Reiterate that this row is a design claim to be "
             "tested, not a measured result.")

# ================================================================== Slide 12
# Roadmap + Close
s = add_slide(NAVY)
add_kicker(s, Inches(0), Inches(0.65), "Roadmap", align=PP_ALIGN.CENTER, width=SLIDE_W)
add_text(s, Inches(1), Inches(1.0), Inches(11.33), Inches(0.8),
          "Proving It", 30, WHITE, bold=True, font=HEADER_FONT, align=PP_ALIGN.CENTER)

rows = [
    "Build natural + controlled occlusion test sets from nuScenes",
    "Fix SORT's circular reference and disabled track-expiry, add ByteTrack, then benchmark all 5 against OATM",
    "Metrics: occluded-object recall, identity preservation, ghost-track rate, recovery time",
    "Split train / validation / test by scene, not by frame, to avoid leakage",
    "Expand from vehicle-only tracking to pedestrians and other vulnerable road users",
]
y = Inches(1.9)
for i, text in enumerate(rows, 1):
    add_badge_row(s, Inches(1.8), y, Inches(9.7), i, text, size=14, row_h=Inches(0.78),
                  text_color=WHITE, badge_color=AMBER_DEEP, number_color=NAVY)
    y += Inches(0.86)

divider = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.67), Inches(6.35), Inches(4), Pt(1.5))
divider.fill.solid()
divider.fill.fore_color.rgb = AMBER
divider.line.fill.background()
divider.shadow.inherit = False
set_fill_alpha(divider, 55)
add_text(s, Inches(1.6), Inches(6.55), Inches(10.13), Inches(0.6),
          "Thank you — questions welcome", 19, AMBER, bold=True, align=PP_ALIGN.CENTER,
          font=HEADER_FONT)
add_notes(s, "Close with the roadmap: this baseline experiment was step one, the frozen-memory "
             "experiment was step two, the SORT test was step three. Next comes fixing the SORT "
             "experiment's methodology, building proper test sets, implementing OATM itself, and "
             "comparing it fairly against existing tracking methods. Also flag scope: everything so far "
             "has been vehicles - pedestrians and other vulnerable road users are the next class to "
             "extend to, since that's arguably the higher-stakes case. Open the floor for questions.")

# ================================================================== Slide 13
# References
s = add_slide(WHITE)
add_kicker(s, MARGIN, Inches(0.55), "References")
add_text(s, MARGIN, Inches(0.95), Inches(11.9), Inches(0.6),
         "Papers We Read", 26, NAVY, bold=True, font=HEADER_FONT)
add_accent_rule(s, MARGIN, Inches(1.62))

references = [
    ("Ma, J., Yang, J., Fu, F., Hao, D., & Wang, R. (2026).",
     "Occlusion-aware multi-camera 3D multi-object tracking with dual-appearance dynamic association.",
     "Applied Soft Computing."),
    ("Pang, Z., Li, J., Tokmakov, P., Chen, D., Zagoruyko, S., & Wang, Y.-X. (2023).",
     "Standing Between Past and Future: Spatio-Temporal Modeling for Multi-Camera 3D Multi-Object Tracking.",
     "CVPR 2023. arXiv:2302.03802."),
    ("Anwar, N., Bilodeau, G.-A., & Bouachir, W. (2024).",
     "STF: Spatio-Temporal Fusion Module for Improving Video Object Detection.",
     "arXiv:2402.10752."),
    ("Cai, J., Xu, M., Li, W., Xiong, Y., Xia, W., Tu, Z., & Soatto, S. (2022).",
     "MeMOT: Multi-Object Tracking with Memory.",
     "CVPR 2022, pp. 8080-8090."),
    ("Tokmakov, P., Li, J., Burgard, W., & Gaidon, A. (2021).",
     "Learning to Track with Object Permanence.",
     "ICCV 2021."),
    ("Geiger, A., Lenz, P., & Urtasun, R. (2012).",
     "Are We Ready for Autonomous Driving? The KITTI Vision Benchmark Suite.",
     "CVPR 2012, pp. 3354-3361."),
    ("Iftikhar, S., Zhang, Z., Asim, M., Muthanna, A., Koucheryavy, A., & Abd El-Latif, A. A. (2022).",
     "Deep Learning-Based Pedestrian Detection in Autonomous Vehicles: Substantial Issues and Challenges.",
     "Electronics, 11(21), 3551."),
]

y = Inches(1.95)
row_h = Inches(0.68)
for authors_year, title, venue in references:
    box = s.shapes.add_textbox(MARGIN, y, CONTENT_W, row_h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = authors_year + " "
    r1.font.size = Pt(11.5)
    r1.font.name = BODY_FONT
    r1.font.color.rgb = DARKGRAY
    r2 = p.add_run()
    r2.text = title + " "
    r2.font.size = Pt(11.5)
    r2.font.italic = True
    r2.font.name = BODY_FONT
    r2.font.color.rgb = NAVY
    r3 = p.add_run()
    r3.text = venue
    r3.font.size = Pt(11.5)
    r3.font.name = BODY_FONT
    r3.font.color.rgb = MUTED
    y += row_h
add_notes(s, "Backup slide for Q&A - full reading list behind the project. One caveat for your own "
             "awareness: the author list for the first entry (the dual-appearance / DADA-Track paper) "
             "came from a secondary search result rather than a directly verified page, since the "
             "publisher page blocked direct access - worth a quick manual check against the actual "
             "journal listing before you present this, just to be safe. Everything else was verified "
             "against the source page or a well-corroborated listing.")

out_path = r"C:\Users\HP\OneDrive\Desktop\srsi\presentation\OATM_Project_Presentation.pptx"
prs.save(out_path)
print("Saved:", out_path)
